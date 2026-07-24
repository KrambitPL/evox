import sys
import pptx
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

def add_sponsor_slide_to_pptx(pptx_path="outputs/evox-simple-explainer.pptx"):
    prs = pptx.Presentation(pptx_path)
    
    if len(prs.slides) >= 6:
        slide = prs.slides[5]
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

    # Colors
    c_navy = RGBColor(15, 23, 42)
    c_slate = RGBColor(71, 85, 105)
    c_blue = RGBColor(37, 99, 235)
    c_green = RGBColor(22, 163, 74)
    c_orange = RGBColor(234, 88, 12)
    c_purple = RGBColor(147, 51, 234)

    # Header
    tb_hdr = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.4))
    tf_hdr = tb_hdr.text_frame
    p_hdr = tf_hdr.paragraphs[0]
    p_hdr.text = "EVOX"
    p_hdr.font.name = "Aptos"
    p_hdr.font.size = Pt(12)
    p_hdr.font.bold = True
    p_hdr.font.color.rgb = c_navy

    tb_page = slide.shapes.add_textbox(Inches(10.58), Inches(0.4), Inches(2.0), Inches(0.4))
    tf_page = tb_page.text_frame
    p_page = tf_page.paragraphs[0]
    p_page.text = "06 / 06"
    p_page.alignment = PP_ALIGN.RIGHT
    p_page.font.name = "Aptos"
    p_page.font.size = Pt(12)
    p_page.font.bold = True
    p_page.font.color.rgb = c_slate

    # Title & Subtitle
    tb_title = slide.shapes.add_textbox(Inches(0.75), Inches(0.85), Inches(11.83), Inches(0.6))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Every sponsor technology is integrated via concrete API adapters."
    p_title.font.name = "Aptos"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = c_navy

    tb_sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(11.83), Inches(0.5))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Real persisted state and configured endpoints power every stage of the Evox governance engine."
    p_sub.font.name = "Aptos"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = c_slate

    # Sponsor cards (5 items with concrete technical implementation details)
    sponsors = [
        (
            "PIONEER",
            "Model Gateway Adapter",
            "evox_api.adapters.pioneer",
            "Wraps the Pioneer REST gateway. Validates model readiness, routes prompt execution, and enforces fail-closed model resolution.",
            c_blue
        ),
        (
            "SENSO",
            "Cited Document Context",
            "evox_api.adapters.senso",
            "Connects to Senso org API to ingest corpus documents and retrieve exact cited text snippets for agent reasoning.",
            c_purple
        ),
        (
            "ACTIAN VECTORAI",
            "Outcome Memory Store",
            "evox_api.adapters.actian",
            "Connects to Actian VectorAI vector DB to index run trajectories and query past failure/success patterns.",
            c_green
        ),
        (
            "BAND",
            "Human Escalation Bridge",
            "evox_api.adapters.band",
            "Uses Band AsyncRestClient to create escalation tasks, poll worker feedback, and deliver human decisions back into jobs.",
            c_orange
        ),
        (
            "GUILD.AI",
            "Release Publisher",
            "evox_api.adapters.guild",
            "Publishes proven candidate workflows to Guild.ai workspace, locking active release versions and evaluator rules.",
            c_navy
        ),
    ]

    card_width = Inches(2.2)
    card_gap = Inches(0.2)
    start_left = Inches(0.75)
    card_top = Inches(2.1)
    card_height = Inches(4.3)

    for i, (name, role, module, desc, color) in enumerate(sponsors):
        left = start_left + i * (card_width + card_gap)
        
        # Background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        # Content inside card
        tb_card = slide.shapes.add_textbox(left + Inches(0.1), card_top + Inches(0.12), card_width - Inches(0.2), card_height - Inches(0.24))
        tf_card = tb_card.text_frame
        tf_card.word_wrap = True
        
        p1 = tf_card.paragraphs[0]
        p1.text = name
        p1.font.name = "Aptos"
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.space_after = Pt(2)

        p2 = tf_card.add_paragraph()
        p2.text = role
        p2.font.name = "Aptos"
        p2.font.size = Pt(10.5)
        p2.font.bold = True
        p2.font.color.rgb = c_navy
        p2.space_after = Pt(2)

        p3 = tf_card.add_paragraph()
        p3.text = module
        p3.font.name = "Courier New"
        p3.font.size = Pt(8.5)
        p3.font.color.rgb = c_slate
        p3.space_after = Pt(8)

        p4 = tf_card.add_paragraph()
        p4.text = desc
        p4.font.name = "Aptos"
        p4.font.size = Pt(9.5)
        p4.font.color.rgb = c_slate

    # Presenter Notes
    slide.notes_slide.notes_text_frame.text = (
        "Concrete Sponsor Integration Architecture:\n"
        "- evox_api.adapters.pioneer: PioneerModelGateway handles model execution and health checks.\n"
        "- evox_api.adapters.senso: SensoAdapter ingests docs and retrieves cited evidence for agent steps.\n"
        "- evox_api.adapters.actian: ActianOutcomeMemorySettings connects to VectorAI for trajectory indexing.\n"
        "- evox_api.adapters.band: EscalationConfig & AsyncRestClient dispatch human approval requests.\n"
        "- evox_api.adapters.guild: GuildSettings verifies agent workspace publication & version governance."
    )

    prs.save(pptx_path)
    print("Updated PPTX with detailed sponsor integration:", pptx_path)

def update_pdf(pdf_path="outputs/evox-simple-explainer.pdf"):
    page_width = 13.33 * 72
    page_height = 7.5 * 72

    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=(page_width, page_height),
        leftMargin=54,
        rightMargin=54,
        topMargin=40,
        bottomMargin=40
    )

    styles = getSampleStyleSheet()

    c_navy = colors.HexColor('#0F172A')
    c_blue = colors.HexColor('#2563EB')
    c_bg_blue = colors.HexColor('#EFF6FF')
    c_green = colors.HexColor('#16A34A')
    c_bg_green = colors.HexColor('#F0FDF4')
    c_orange = colors.HexColor('#EA580C')
    c_bg_orange = colors.HexColor('#FFF7ED')
    c_slate = colors.HexColor('#475569')
    c_light_gray = colors.HexColor('#E2E8F0')
    c_purple = colors.HexColor('#9333EA')
    c_bg_purple = colors.HexColor('#FAF5FF')

    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=26,
        leading=32,
        textColor=c_navy
    )

    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=15,
        leading=20,
        textColor=c_slate
    )

    header_style = ParagraphStyle(
        'HeaderStyle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=14,
        textColor=c_navy
    )

    header_right_style = ParagraphStyle(
        'HeaderRightStyle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=14,
        textColor=c_slate,
        alignment=2
    )

    box_header_style = ParagraphStyle(
        'BoxHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=c_navy,
        alignment=1
    )

    box_body_style = ParagraphStyle(
        'BoxBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=13.5,
        textColor=c_slate,
        alignment=1
    )

    code_style = ParagraphStyle(
        'CodeText',
        parent=styles['Normal'],
        fontName='Courier-Oblique',
        fontSize=8.5,
        leading=11,
        textColor=c_slate,
        alignment=1
    )

    row_label_style = ParagraphStyle(
        'RowLabel',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=c_navy
    )

    row_body_style = ParagraphStyle(
        'RowBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=13,
        leading=16,
        textColor=c_slate
    )

    bottom_note_style = ParagraphStyle(
        'BottomNote',
        parent=styles['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=12.5,
        leading=16,
        textColor=c_slate,
        alignment=1
    )

    story = []

    def make_header(slide_num):
        t = Table([
            [Paragraph("EVOX", header_style), Paragraph(f"0{slide_num} / 06", header_right_style)]
        ], colWidths=[400, 452])
        t.setStyle(TableStyle([
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOTTOMPADDING', (0,0), (-1,-1), 0),
            ('TOPPADDING', (0,0), (-1,-1), 0),
        ]))
        return t

    # ---------------- SLIDE 1 ----------------
    story.append(make_header(1))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=25))

    story.append(Paragraph("Evox takes a problem, builds an agentic system, and improves it.", title_style))
    story.append(Spacer(1, 12))
    story.append(Paragraph("You define what success means. Evox turns that definition into a system that can learn from every measured result.", subtitle_style))
    story.append(Spacer(1, 45))

    b1 = Paragraph("<b>PROBLEM</b>", box_header_style)
    b2 = Paragraph("<b>AGENTIC SYSTEM</b>", ParagraphStyle('B2', parent=box_header_style, textColor=c_blue))
    b3 = Paragraph("<b>BETTER RESULTS</b>", ParagraphStyle('B3', parent=box_header_style, textColor=c_green))
    arr = Paragraph("<b>→</b>", ParagraphStyle('Arr', parent=box_header_style, fontSize=24, textColor=c_slate))

    flow_table = Table([[b1, arr, b2, arr, b3]], colWidths=[230, 40, 230, 40, 230])
    flow_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), c_light_gray),
        ('BACKGROUND', (2,0), (2,0), c_bg_blue),
        ('BACKGROUND', (4,0), (4,0), c_bg_green),
        ('BOX', (0,0), (0,0), 1, c_slate),
        ('BOX', (2,0), (2,0), 1, c_blue),
        ('BOX', (4,0), (4,0), 1, c_green),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 25),
        ('BOTTOMPADDING', (0,0), (-1,-1), 25),
    ]))
    story.append(flow_table)

    # ---------------- SLIDE 2 ----------------
    story.append(PageBreak())
    story.append(make_header(2))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=20))

    story.append(Paragraph("1 &nbsp; Start with the problem—and define what success means.", title_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Evox needs four things before it can build honestly.", subtitle_style))
    story.append(Spacer(1, 25))

    rows_data = [
        [Paragraph("PROBLEM", row_label_style), Paragraph("What needs to be solved?", row_body_style)],
        [Paragraph("SUCCESS", ParagraphStyle('R2', parent=row_label_style, textColor=c_green)), Paragraph("What result counts as good?", row_body_style)],
        [Paragraph("EVIDENCE", row_label_style), Paragraph("What data can prove the result?", row_body_style)],
        [Paragraph("BOUNDARIES", ParagraphStyle('R4', parent=row_label_style, textColor=c_navy)), Paragraph("What may the system do—and when must it ask a human?", row_body_style)],
    ]
    grid_table = Table(rows_data, colWidths=[180, 672])
    grid_table.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 14),
        ('BOTTOMPADDING', (0,0), (-1,-1), 14),
        ('LINEBELOW', (0,0), (-1,-1), 1, c_light_gray),
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#FAFAFA')),
        ('LEFTPADDING', (0,0), (-1,-1), 16),
    ]))
    story.append(grid_table)
    story.append(Spacer(1, 25))
    story.append(Paragraph("If success cannot be measured, the system cannot improve honestly.", bottom_note_style))

    # ---------------- SLIDE 3 ----------------
    story.append(PageBreak())
    story.append(make_header(3))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=20))

    story.append(Paragraph("2 &nbsp; Evox turns that mission into an agentic system.", title_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("The system is designed around the job—not around a fixed business template.", subtitle_style))
    story.append(Spacer(1, 35))

    steps_data = [
        [
            Paragraph("<b>UNDERSTAND</b>", box_header_style),
            Paragraph("<b>→</b>", ParagraphStyle('A1', parent=box_header_style, fontSize=20, textColor=c_slate)),
            Paragraph("<b>PLAN</b>", box_header_style),
            Paragraph("<b>→</b>", ParagraphStyle('A2', parent=box_header_style, fontSize=20, textColor=c_slate)),
            Paragraph("<b>ACT</b>", box_header_style),
            Paragraph("<b>→</b>", ParagraphStyle('A3', parent=box_header_style, fontSize=20, textColor=c_slate)),
            Paragraph("<b>CHECK</b>", box_header_style),
        ],
        [
            Paragraph("Read the problem and available evidence.", box_body_style),
            "",
            Paragraph("Break the mission into agent steps.", box_body_style),
            "",
            Paragraph("Use the permitted models, tools, and data.", box_body_style),
            "",
            Paragraph("Verify the result or escalate to a human.", box_body_style),
        ]
    ]
    step_table = Table(steps_data, colWidths=[180, 30, 180, 30, 180, 30, 180])
    step_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,1), c_bg_blue),
        ('BACKGROUND', (2,0), (2,1), c_bg_blue),
        ('BACKGROUND', (4,0), (4,1), c_bg_blue),
        ('BACKGROUND', (6,0), (6,1), c_bg_blue),
        ('BOX', (0,0), (0,1), 1, c_blue),
        ('BOX', (2,0), (2,1), 1, c_blue),
        ('BOX', (4,0), (4,1), 1, c_blue),
        ('BOX', (6,0), (6,1), 1, c_blue),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,0), 16),
        ('BOTTOMPADDING', (0,1), (-1,1), 16),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('RIGHTPADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(step_table)
    story.append(Spacer(1, 45))
    story.append(Paragraph("The result is a versioned workflow that can be tested, compared, and improved.", bottom_note_style))

    # ---------------- SLIDE 4 ----------------
    story.append(PageBreak())
    story.append(make_header(4))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=20))

    story.append(Paragraph("3 &nbsp; Every result is judged against the same success criteria.", title_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("That fixed definition of success becomes the learning signal.", subtitle_style))
    story.append(Spacer(1, 30))

    c_box_hdr = Paragraph("<b>SUCCESS<br/>CRITERIA</b>", ParagraphStyle('CBH', parent=box_header_style, textColor=colors.white))
    c_box_bdy = Paragraph("The target stays fixed while candidates change.", ParagraphStyle('CBB', parent=box_body_style, textColor=colors.HexColor('#F8FAFC')))

    s_hdr = Paragraph("<b>SUCCESS</b>", ParagraphStyle('SH', parent=box_header_style, textColor=c_green))
    s_bdy = Paragraph("Reinforce the steps that worked.", box_body_style)
    f_hdr = Paragraph("<b>FAILURE</b>", ParagraphStyle('FH', parent=box_header_style, textColor=c_orange))
    f_bdy = Paragraph("Find the weak step and create a better candidate.", box_body_style)
    u_hdr = Paragraph("<b>UNCERTAIN</b>", ParagraphStyle('UH', parent=box_header_style, textColor=c_navy))
    u_bdy = Paragraph("Ask a human instead of guessing.", box_body_style)

    arrow_large = Paragraph("<b>→</b>", ParagraphStyle('AL', parent=box_header_style, fontSize=28, textColor=c_slate))

    outcomes_table = Table([[s_hdr, s_bdy], [f_hdr, f_bdy], [u_hdr, u_bdy]], colWidths=[140, 360])
    outcomes_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (1,0), c_bg_green),
        ('BOX', (0,0), (1,0), 1, c_green),
        ('BACKGROUND', (0,1), (1,1), c_bg_orange),
        ('BOX', (0,1), (1,1), 1, c_orange),
        ('BACKGROUND', (0,2), (1,2), colors.HexColor('#F8FAFC')),
        ('BOX', (0,2), (1,2), 1, c_slate),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 12),
        ('BOTTOMPADDING', (0,0), (-1,-1), 12),
        ('LEFTPADDING', (0,0), (-1,-1), 12),
    ]))

    left_box_table = Table([[c_box_hdr], [c_box_bdy]], colWidths=[240])
    left_box_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,1), c_navy),
        ('BOX', (0,0), (0,1), 1, c_navy),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,0), 24),
        ('BOTTOMPADDING', (0,1), (-1,1), 24),
        ('LEFTPADDING', (0,0), (-1,-1), 12),
        ('RIGHTPADDING', (0,0), (-1,-1), 12),
    ]))

    eval_layout = Table([[left_box_table, arrow_large, outcomes_table]], colWidths=[250, 40, 510])
    eval_layout.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('ALIGN', (1,0), (1,0), 'CENTER'),
    ]))
    story.append(eval_layout)
    story.append(Spacer(1, 30))
    story.append(Paragraph("Successes and failures become evidence for the next version.", bottom_note_style))

    # ---------------- SLIDE 5 ----------------
    story.append(PageBreak())
    story.append(make_header(5))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=20))

    story.append(Paragraph("4 &nbsp; The next version must prove it is better.", title_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Evox improves the workflow without letting the workflow rewrite the rules.", subtitle_style))
    story.append(Spacer(1, 25))

    p1_hdr = Paragraph("<b>RUN</b>", box_header_style)
    p1_bdy = Paragraph("Use the current approved system.", box_body_style)
    p2_hdr = Paragraph("<b>LEARN</b>", box_header_style)
    p2_bdy = Paragraph("Study measured successes and failures.", box_body_style)
    p3_hdr = Paragraph("<b>CANDIDATE</b>", box_header_style)
    p3_bdy = Paragraph("Create a new version of the workflow.", box_body_style)
    p4_hdr = Paragraph("<b>PROVE</b>", ParagraphStyle('P4H', parent=box_header_style, textColor=c_green))
    p4_bdy = Paragraph("Promote only if it performs better and stays allowed.", box_body_style)

    loop_table = Table([
        [p1_hdr, arr, p2_hdr, arr, p3_hdr, arr, p4_hdr],
        [p1_bdy, "", p2_bdy, "", p3_bdy, "", p4_bdy]
    ], colWidths=[180, 30, 180, 30, 180, 30, 180])
    loop_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,1), colors.HexColor('#F8FAFC')),
        ('BACKGROUND', (2,0), (2,1), colors.HexColor('#F8FAFC')),
        ('BACKGROUND', (4,0), (4,1), colors.HexColor('#F8FAFC')),
        ('BACKGROUND', (6,0), (6,1), c_bg_green),
        ('BOX', (0,0), (0,1), 1, c_slate),
        ('BOX', (2,0), (2,1), 1, c_slate),
        ('BOX', (4,0), (4,1), 1, c_slate),
        ('BOX', (6,0), (6,1), 1, c_green),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,0), 12),
        ('BOTTOMPADDING', (0,1), (-1,1), 12),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]))
    story.append(loop_table)
    story.append(Spacer(1, 20))

    dec_left = Paragraph("<b>IF IT FAILS:</b> KEEP THE CURRENT VERSION", ParagraphStyle('DL', parent=box_body_style, textColor=c_orange, fontSize=12))
    dec_right = Paragraph("<b>IF IT WINS:</b> MAKE IT THE NEW VERSION", ParagraphStyle('DR', parent=box_body_style, textColor=c_green, fontSize=12))

    dec_table = Table([[dec_left, dec_right]], colWidths=[400, 400])
    dec_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), c_bg_orange),
        ('BOX', (0,0), (0,0), 1, c_orange),
        ('BACKGROUND', (1,0), (1,0), c_bg_green),
        ('BOX', (1,0), (1,0), 1, c_green),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 10),
        ('BOTTOMPADDING', (0,0), (-1,-1), 10),
    ]))
    story.append(dec_table)
    story.append(Spacer(1, 25))

    guardrail_box = Table([[Paragraph("<b>The system learns how to work. You keep control of what “better” means.</b>", ParagraphStyle('GB', parent=box_header_style, textColor=c_navy, fontSize=13))]], colWidths=[852])
    guardrail_box.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), c_light_gray),
        ('BOX', (0,0), (0,0), 1, c_slate),
        ('VALIGN', (0,0), (0,0), 'MIDDLE'),
        ('TOPPADDING', (0,0), (0,0), 12),
        ('BOTTOMPADDING', (0,0), (0,0), 12),
    ]))
    story.append(guardrail_box)

    # ---------------- SLIDE 6 (CONCRETE SPONSOR INTEGRATION DETAILS) ----------------
    story.append(PageBreak())
    story.append(make_header(6))
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=1, color=c_light_gray, spaceBefore=5, spaceAfter=18))

    story.append(Paragraph("Every sponsor technology is integrated via concrete API adapters.", title_style))
    story.append(Spacer(1, 8))
    story.append(Paragraph("Real persisted state and configured endpoints power every stage of the Evox governance engine.", subtitle_style))
    story.append(Spacer(1, 20))

    card_pioneer = [
        Paragraph("<b>PIONEER</b>", ParagraphStyle('CP1', parent=box_header_style, textColor=c_blue)),
        Paragraph("Model Gateway Adapter", ParagraphStyle('CR1', parent=box_header_style, textColor=c_navy, fontSize=11)),
        Paragraph("evox_api.adapters.pioneer", code_style),
        Spacer(1, 6),
        Paragraph("Wraps Pioneer REST gateway. Validates model readiness, routes prompt execution, and enforces fail-closed model resolution.", box_body_style)
    ]

    card_senso = [
        Paragraph("<b>SENSO</b>", ParagraphStyle('CP2', parent=box_header_style, textColor=c_purple)),
        Paragraph("Cited Context Adapter", ParagraphStyle('CR2', parent=box_header_style, textColor=c_navy, fontSize=11)),
        Paragraph("evox_api.adapters.senso", code_style),
        Spacer(1, 6),
        Paragraph("Connects to Senso org API to ingest corpus documents and retrieve exact cited text snippets for agent reasoning.", box_body_style)
    ]

    card_actian = [
        Paragraph("<b>ACTIAN VECTORAI</b>", ParagraphStyle('CP3', parent=box_header_style, textColor=c_green)),
        Paragraph("Outcome Memory Store", ParagraphStyle('CR3', parent=box_header_style, textColor=c_navy, fontSize=11)),
        Paragraph("evox_api.adapters.actian", code_style),
        Spacer(1, 6),
        Paragraph("Connects to Actian VectorAI vector DB to index run trajectories and query past failure/success patterns.", box_body_style)
    ]

    card_band = [
        Paragraph("<b>BAND</b>", ParagraphStyle('CP4', parent=box_header_style, textColor=c_orange)),
        Paragraph("Human Escalation Bridge", ParagraphStyle('CR4', parent=box_header_style, textColor=c_navy, fontSize=11)),
        Paragraph("evox_api.adapters.band", code_style),
        Spacer(1, 6),
        Paragraph("Uses Band AsyncRestClient to create escalation tasks, poll worker feedback, and deliver human decisions back into jobs.", box_body_style)
    ]

    card_guild = [
        Paragraph("<b>GUILD.AI</b>", ParagraphStyle('CP5', parent=box_header_style, textColor=c_navy)),
        Paragraph("Release Publisher", ParagraphStyle('CR5', parent=box_header_style, textColor=c_navy, fontSize=11)),
        Paragraph("evox_api.adapters.guild", code_style),
        Spacer(1, 6),
        Paragraph("Publishes proven candidate workflows to Guild.ai workspace, locking active release versions and evaluator rules.", box_body_style)
    ]

    sponsors_grid = Table([[card_pioneer, card_senso, card_actian, card_band, card_guild]], colWidths=[162, 162, 162, 162, 162])
    sponsors_grid.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), c_bg_blue),
        ('BOX', (0,0), (0,0), 1, c_blue),
        ('BACKGROUND', (1,0), (1,0), c_bg_purple),
        ('BOX', (1,0), (1,0), 1, c_purple),
        ('BACKGROUND', (2,0), (2,0), c_bg_green),
        ('BOX', (2,0), (2,0), 1, c_green),
        ('BACKGROUND', (3,0), (3,0), c_bg_orange),
        ('BOX', (3,0), (3,0), 1, c_orange),
        ('BACKGROUND', (4,0), (4,0), colors.HexColor('#F8FAFC')),
        ('BOX', (4,0), (4,0), 1, c_navy),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 14),
        ('BOTTOMPADDING', (0,0), (-1,-1), 14),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]))

    story.append(sponsors_grid)
    story.append(Spacer(1, 25))
    story.append(Paragraph("Integrations are fail-closed with real persisted state — no synthetic or mock fallbacks in production paths.", bottom_note_style))

    doc.build(story)
    print("Updated PDF:", pdf_path)

if __name__ == "__main__":
    add_sponsor_slide_to_pptx()
    update_pdf()
